#!/usr/bin/env python3
"""
srt_extract_slides.py — Extract terminology from video slide screenshots.

Pipeline: ffmpeg frame extraction → perceptual hash dedup → gemma4:26b VLM (Ollama) → terms file

Usage:
    python3 srt_extract_slides.py <video_file> [--output <terms_file>] [--interval 60] [--model <model_name>]

Output:
    Writes _slide_terms.txt in the same directory as the video file (or --output path).
    Format is compatible with srt skill's existing terminology injection.
"""

from __future__ import annotations

import argparse
import importlib
import json
import os
import platform
import re
import shutil
import subprocess
import sys
import tempfile
import time

OLLAMA_DEFAULT_MODEL = "gemma4:26b"
MLX_DEFAULT_MODEL = "lmstudio-community/Qwen3-VL-8B-Instruct-MLX-4bit"
OCR_ENGINES = {"apple-vision", "rapidocr"}
UI_CHROME_STOPLIST = {
    "OBS",
    "OBS Studio",
    "Pointofix",
    "Microsoft PowerPoint",
    "PowerPoint",
    "iSlide",
    "Tiger Trade",
    "自選",
    "個股資料",
}


def extract_frames(video_path: str, output_dir: str, interval: int = 60) -> list[tuple[str, float]]:
    """Extract one frame per interval seconds using ffmpeg.
    Returns list of (path, time_seconds) tuples."""
    if shutil.which("ffmpeg") is None:
        raise RuntimeError(
            "ffmpeg not found. Install ffmpeg first: macOS `brew install ffmpeg`; "
            "Windows `winget install Gyan.FFmpeg` or `choco install ffmpeg`; "
            "Linux `sudo apt install ffmpeg`."
        )
    pattern = os.path.join(output_dir, "frame_%04d.jpg")
    cmd = [
        "ffmpeg", "-i", video_path,
        "-vf", f"fps=1/{interval}",
        "-q:v", "5",
        pattern, "-y"
    ]
    subprocess.run(cmd, capture_output=True, check=True)
    frames = sorted(
        [f for f in os.listdir(output_dir) if f.startswith("frame_") and f.endswith(".jpg")]
    )
    # frame_0001.jpg = interval seconds, frame_0002.jpg = 2*interval, etc.
    return [(os.path.join(output_dir, f), (i + 1) * interval) for i, f in enumerate(frames)]


def deduplicate_frames(frames: list[tuple[str, float]], threshold: int = 8) -> list[tuple[str, float]]:
    """Remove visually similar consecutive frames using perceptual hashing.
    Input/output: list of (path, time_seconds) tuples."""
    import imagehash
    from PIL import Image

    unique = []
    seen_hashes = []

    for frame_path, time_s in frames:
        img = Image.open(frame_path)
        h = imagehash.phash(img)

        is_dup = any(h - prev_h < threshold for prev_h in seen_hashes)
        if not is_dup:
            seen_hashes.append(h)
            unique.append((frame_path, time_s))

    return unique


def _get_vlm_prompt(caption_mode: bool) -> str:
    if caption_mode:
        return """用 2-3 句話描述這張畫面的內容，包含：
- 投影片標題或主題
- 畫面上可見的所有英文術語、ticker、人名（保留原始拼法）
- 如果是圖表，描述圖表類型和主要趨勢

輸出 JSON：{"caption": "...", "terms": ["term1", "term2"]}
只輸出 JSON，不要其他文字。"""
    return """請仔細看這張投影片截圖，抽取所有可見的文字資訊，特別是：
1. 股票代號 (ticker symbols)，例如 AAPL, TSLA
2. 專有名詞（公司名、人名、品牌名）
3. 中文和英文的金融/技術術語
4. 重要的數字和日期

用 JSON 格式輸出：{"tickers": [...], "proper_nouns": [...], "technical_terms": [...], "slide_title": "..."}
只輸出 JSON，不要其他文字。"""


def ocr_with_ollama_vlm(frames: list[tuple[str, float]], model_name: str, caption_mode: bool = False) -> list[dict]:
    """Run Ollama vision model on each frame. Works with gemma4, llava, etc."""
    import base64

    requests = importlib.import_module("requests")
    prompt = _get_vlm_prompt(caption_mode)
    url = "http://localhost:11434/api/chat"

    results = []
    for i, (frame_path, frame_time) in enumerate(frames):
        with open(frame_path, "rb") as f:
            img_b64 = base64.b64encode(f.read()).decode()

        t0 = time.time()
        resp = requests.post(url, json={
            "model": model_name,
            "messages": [{"role": "user", "content": prompt, "images": [img_b64]}],
            "stream": False,
            "options": {"temperature": 0.1, "num_predict": 512},
            "think": False,
        }, timeout=300)
        elapsed = time.time() - t0

        resp.raise_for_status()
        data = resp.json()
        if isinstance(data, dict) and data.get("error"):
            raise RuntimeError(f"ollama error: {data['error']}")
        output_text = data.get("message", {}).get("content", "")
        print(f"  [{i+1}/{len(frames)}] {os.path.basename(frame_path)} (t={frame_time:.0f}s) — {elapsed:.1f}s", file=sys.stderr)

        results.append({"frame": frame_path, "frame_time": frame_time, "inference_time": elapsed, "raw": output_text})

    return results


def ocr_with_qwen_vlm(frames: list[tuple[str, float]], model_name: str, caption_mode: bool = False) -> list[dict]:
    """Run Qwen VL on each frame to extract terminology or captions.
    Input: list of (path, time_seconds) tuples."""
    from mlx_vlm import load, generate
    from mlx_vlm.prompt_utils import apply_chat_template
    from mlx_vlm.utils import load_config

    model, processor = load(model_name)
    config = load_config(model_name)
    prompt = _get_vlm_prompt(caption_mode)

    results = []
    for i, (frame_path, frame_time) in enumerate(frames):
        t0 = time.time()
        formatted = apply_chat_template(processor, config, prompt, num_images=1)
        result = generate(model, processor, formatted, [frame_path], max_tokens=512, verbose=False)
        elapsed = time.time() - t0

        output_text = result.text if hasattr(result, "text") else str(result)
        print(f"  [{i+1}/{len(frames)}] {os.path.basename(frame_path)} (t={frame_time:.0f}s) — {elapsed:.1f}s", file=sys.stderr)

        results.append({"frame": frame_path, "frame_time": frame_time, "inference_time": elapsed, "raw": output_text})

    return results


def _apple_vision_available() -> bool:
    try:
        import Vision
        import Quartz  # noqa: F401
    except ImportError:
        return False
    return hasattr(Vision, "VNRecognizeTextRequest")


def _rapidocr_importable() -> bool:
    try:
        importlib.import_module("rapidocr")
    except ImportError:
        return False
    return True


def _warn(message: str):
    print(f"  WARNING: {message}", file=sys.stderr)


def resolve_engine(args) -> tuple[str, str | None]:
    engine = args.engine
    model = args.model

    if engine == "auto":
        if model is not None:
            return ("mlx", model) if "/" in model else ("ollama", model)
        if _rapidocr_importable():
            return "rapidocr", None
        if platform.system() == "Darwin" and _apple_vision_available():
            return "apple-vision", None
        raise RuntimeError('RapidOCR is required for auto OCR. Install with: pip install "rapidocr>=3.9,<4" onnxruntime')

    if engine == "rapidocr":
        if model is not None:
            _warn("rapidocr ignores --model")
        if not _rapidocr_importable():
            raise RuntimeError('rapidocr is not installed. Install with: pip install "rapidocr>=3.9,<4" onnxruntime')
        return "rapidocr", None

    if engine == "apple-vision":
        if model is not None:
            _warn("apple-vision ignores --model")
        if platform.system() != "Darwin" or not _apple_vision_available():
            raise RuntimeError("apple-vision is only available on macOS with pyobjc Vision/Quartz")
        return "apple-vision", None

    if engine == "ollama":
        if model is None:
            return "ollama", OLLAMA_DEFAULT_MODEL
        return "ollama", model

    if engine == "mlx":
        if model is None:
            return "mlx", MLX_DEFAULT_MODEL
        return "mlx", model

    raise RuntimeError(f"unknown engine: {engine}")


def _recognize_text_with_vision(frame_path: str, request) -> str:
    import Quartz
    from Foundation import NSURL
    import Vision

    url = NSURL.fileURLWithPath_(frame_path)
    source = Quartz.CGImageSourceCreateWithURL(url, None)
    image = Quartz.CGImageSourceCreateImageAtIndex(source, 0, None) if source else None
    if image is None:
        raise RuntimeError(f"failed to load image: {frame_path}")

    handler = Vision.VNImageRequestHandler.alloc().initWithCGImage_options_(image, {})
    ok = handler.performRequests_error_([request], None)
    if isinstance(ok, tuple) and ok and not ok[0]:
        raise RuntimeError(ok[1] or "Vision request failed")
    if ok is False:
        raise RuntimeError("Vision request failed")

    lines = []
    for observation in request.results() or []:
        candidates = observation.topCandidates_(1)
        if candidates:
            lines.append(str(candidates[0].string()))
    return "\n".join(lines)


def _ocr_with_apple_vision_languages(frames: list[tuple[str, float]], languages: tuple[str, ...]) -> list[dict]:
    import Vision

    request = Vision.VNRecognizeTextRequest.alloc().init()
    request.setRecognitionLevel_(Vision.VNRequestTextRecognitionLevelAccurate)
    request.setUsesLanguageCorrection_(False)
    request.setRecognitionLanguages_(list(languages))
    request.setMinimumTextHeight_(0.0)

    results = []
    failed = 0
    last_error = None
    for i, (frame_path, frame_time) in enumerate(frames):
        t0 = time.time()
        try:
            raw = _recognize_text_with_vision(frame_path, request)
        except Exception as exc:
            failed += 1
            last_error = exc
            _warn(f"Apple Vision OCR failed for {frame_path}: {exc}")
            raw = ""
        elapsed = time.time() - t0
        print(f"  [{i+1}/{len(frames)}] {os.path.basename(frame_path)} (t={frame_time:.0f}s) — {elapsed:.1f}s", file=sys.stderr)
        results.append({"frame": frame_path, "frame_time": frame_time, "inference_time": elapsed, "raw": raw})
    if frames and failed == len(frames):
        raise RuntimeError(f"Apple Vision OCR failed for all frames: {last_error}")
    return results


def ocr_with_apple_vision(frames: list[tuple[str, float]], languages: tuple[str, ...] = ("zh-Hant", "en-US")) -> list[dict]:
    """Run macOS Apple Vision OCR on each frame."""
    try:
        return _ocr_with_apple_vision_languages(frames, tuple(languages))
    except Exception as exc:
        if tuple(languages) != ("en-US",):
            _warn(f"Apple Vision OCR failed with {languages}; retrying en-US only: {exc}")
            return _ocr_with_apple_vision_languages(frames, ("en-US",))
        raise


def _rapidocr_result_lines(result) -> list[str]:
    txts = getattr(result, "txts", None)
    if txts is not None:
        return [str(text).strip() for text in txts if str(text).strip()]

    to_json = getattr(result, "to_json", None)
    if callable(to_json):
        data = to_json() or []
        return [str(item.get("text", item.get("txt", ""))).strip() for item in data if str(item.get("text", item.get("txt", ""))).strip()]

    lines = []
    for item in result or []:
        if isinstance(item, (list, tuple)) and len(item) >= 2:
            text = item[1]
            if isinstance(text, (list, tuple)) and text:
                text = text[0]
            text = str(text).strip()
            if text:
                lines.append(text)
    return lines


def ocr_with_rapidocr(frames: list[tuple[str, float]], lang: str | None = None) -> list[dict]:
    """Run RapidOCR v3 on each frame and return plain OCR text.

    Uses the default PP-OCRv5 ``ch`` model, which empirically beats the dedicated
    ``chinese_cht`` v3 model on Fred's mixed English-ticker + Traditional-Chinese
    screens (11 vs 1 ticker hits on a real Tiger Trade frame, 2026-06-28): the v5
    ``ch`` model reads English tickers cleanly AND Traditional Chinese correctly,
    while the older chinese_cht v3 rec model mangles English. ``lang`` is reserved
    for an explicit override and is currently unused.
    """
    try:
        from rapidocr import RapidOCR
    except ImportError as exc:
        raise RuntimeError('rapidocr is not installed. Install with: pip install "rapidocr>=3.9,<4" onnxruntime') from exc

    try:
        engine = RapidOCR()
    except Exception as exc:
        raise RuntimeError(
            "RapidOCR engine construction failed. Likely causes: (1) the ONNX Runtime backend is "
            "not installed (pip install onnxruntime), or (2) the OCR model is missing or its cache "
            "is unavailable (allow network on first run, or pre-load the model into the RapidOCR "
            f"cache for offline/Docker). Original error: {exc}"
        ) from exc

    results = []
    failed = 0
    last_error = None
    for i, (frame_path, frame_time) in enumerate(frames):
        t0 = time.time()
        try:
            # OCR output is text-only by design; bbox/score intentionally not persisted in this version.
            raw = "\n".join(_rapidocr_result_lines(engine(frame_path)))
        except Exception as exc:
            failed += 1
            last_error = exc
            _warn(f"RapidOCR failed for {frame_path}: {exc}")
            raw = ""
        elapsed = time.time() - t0
        print(f"  [{i+1}/{len(frames)}] {os.path.basename(frame_path)} (t={frame_time:.0f}s) — {elapsed:.1f}s", file=sys.stderr)
        results.append({"frame": frame_path, "frame_time": frame_time, "inference_time": elapsed, "raw": raw})
    if frames and failed == len(frames):
        raise RuntimeError(f"RapidOCR failed for all frames: {last_error}")
    return results


def parse_vlm_outputs(results: list[dict], caption_mode: bool = False) -> dict:
    """Parse VLM JSON outputs. In caption mode, returns timestamped captions.
    In legacy mode, merges into a unified terminology set."""
    if caption_mode:
        captions = []
        all_terms = set()
        for r in results:
            raw = r["raw"]
            json_match = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", raw, re.DOTALL)
            if json_match:
                raw = json_match.group(1)
            try:
                data = json.loads(raw)
                caption = data.get("caption", "")
                terms = data.get("terms", [])
                if caption:
                    captions.append({
                        "time_s": r["frame_time"],
                        "caption": caption,
                        "terms": terms,
                    })
                    all_terms.update(terms)
            except json.JSONDecodeError:
                print(f"  WARNING: Failed to parse JSON from {r['frame']}", file=sys.stderr)
        return {"captions": captions, "all_terms": sorted(all_terms)}

    # Legacy OCR mode
    all_tickers = set()
    all_proper_nouns = set()
    all_terms = set()
    all_titles = set()

    for r in results:
        raw = r["raw"]
        json_match = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", raw, re.DOTALL)
        if json_match:
            raw = json_match.group(1)
        try:
            data = json.loads(raw)
            all_tickers.update(data.get("tickers", []))
            all_proper_nouns.update(data.get("proper_nouns", []))
            all_terms.update(data.get("technical_terms", []))
            title = data.get("slide_title", "")
            if title:
                all_titles.add(title)
        except json.JSONDecodeError:
            print(f"  WARNING: Failed to parse JSON from {r['frame']}", file=sys.stderr)

    return {
        "tickers": sorted(all_tickers),
        "proper_nouns": sorted(all_proper_nouns),
        "technical_terms": sorted(all_terms),
        "slide_titles": sorted(all_titles),
    }


def _ocr_lines(results: list[dict]) -> list[str]:
    lines = []
    seen = set()
    for r in results:
        for line in r.get("raw", "").splitlines():
            line = line.strip()
            if not line or line in seen:
                continue
            seen.add(line)
            lines.append(line)
    return lines


def _conservative_tickers(lines: list[str]) -> list[str]:
    tickers = set()
    for line in lines:
        for token in re.findall(r"\b[A-Z]{2,5}\b", line):
            if token not in UI_CHROME_STOPLIST:
                tickers.add(token)
    return sorted(tickers)


def parse_ocr_outputs(results: list[dict], caption_mode: bool = False) -> dict:
    """Parse plain OCR text without assuming VLM JSON."""
    lines = _ocr_lines(results)
    tickers = _conservative_tickers(lines)

    if caption_mode:
        captions = []
        for r in results:
            raw = r.get("raw", "").strip()
            if raw:
                captions.append({"time_s": r["frame_time"], "caption": raw, "terms": []})
        return {"captions": captions, "all_terms": tickers}

    return {
        "tickers": tickers,
        "proper_nouns": [],
        "technical_terms": [],
        "slide_titles": [],
        "raw_ocr": lines,
    }


def run_vlm_engine(engine: str, frames: list[tuple[str, float]], model: str, caption_mode: bool) -> list[dict]:
    if engine == "ollama":
        return ocr_with_ollama_vlm(frames, model, caption_mode=caption_mode)
    return ocr_with_qwen_vlm(frames, model, caption_mode=caption_mode)


def write_terms_file(terms: dict, output_path: str):
    """Write terminology file compatible with srt skill's _slide_terms.txt format."""
    lines = []
    lines.append("# 本集投影片自動抽取的術語")
    lines.append(f"# 抽取時間: {time.strftime('%Y-%m-%d %H:%M')}")
    lines.append("")

    if terms.get("tickers"):
        lines.append("# Ticker Symbols")
        for t in terms["tickers"]:
            lines.append(t)
        lines.append("")

    if terms.get("proper_nouns"):
        lines.append("# 專有名詞")
        for n in terms["proper_nouns"]:
            lines.append(n)
        lines.append("")

    if terms.get("technical_terms"):
        lines.append("# 技術/金融術語")
        for t in terms["technical_terms"]:
            lines.append(t)
        lines.append("")

    if terms.get("slide_titles"):
        lines.append("# 投影片標題")
        for t in terms["slide_titles"]:
            lines.append(t)
        lines.append("")

    if terms.get("raw_ocr"):
        lines.append("# 螢幕 OCR 文字（原始）")
        for line in terms["raw_ocr"]:
            lines.append(line)
        lines.append("")

    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))


def extract_pptx_text(pptx_path: str) -> list[str]:
    """Extract text lines from PowerPoint slides, tables, and speaker notes."""
    try:
        from pptx import Presentation
    except ImportError:
        print("ERROR: python-pptx is required for .pptx/.ppt input. Install with: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(pptx_path)
    lines = []

    def add_line(text: str):
        text = text.strip()
        if text:
            lines.append(text)

    def walk_shapes(shapes):
        for sh in shapes:
            if getattr(sh, "shape_type", None) == 6 and hasattr(sh, "shapes"):
                walk_shapes(sh.shapes)
                continue

            if getattr(sh, "has_table", False):
                for row in sh.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    row_text = " | ".join(cells).strip()
                    if row_text.replace("|", "").strip():
                        add_line(row_text)

            if getattr(sh, "has_text_frame", False):
                for para in sh.text_frame.paragraphs:
                    add_line("".join(run.text for run in para.runs))

    for slide in prs.slides:
        walk_shapes(slide.shapes)
        if slide.has_notes_slide:
            add_line(slide.notes_slide.notes_text_frame.text)

    unique_lines = []
    seen = set()
    for line in lines:
        if line not in seen:
            seen.add(line)
            unique_lines.append(line)

    return unique_lines


def main():
    parser = argparse.ArgumentParser(description="Extract slide terminology from video")
    parser.add_argument("video", help="Path to video file or .pptx/.ppt slide file")
    parser.add_argument("--output", "-o", help="Output terms file path (default: <video_dir>/_slide_terms.txt)")
    parser.add_argument("--interval", type=int, default=60, help="Frame extraction interval in seconds (default: 60)")
    parser.add_argument("--threshold", type=int, default=8, help="Perceptual hash dedup threshold (default: 8)")
    parser.add_argument("--engine", choices=("auto", "rapidocr", "apple-vision", "ollama", "mlx"), default="auto",
                        help="OCR engine (default: auto)")
    parser.add_argument("--model", default=None,
                        help="VLM model name (auto infers mlx when the model contains '/', otherwise ollama)")
    parser.add_argument("--json", action="store_true", help="Also output raw JSON results")
    parser.add_argument("--caption", action="store_true",
                        help="Caption mode: output timestamped captions + terms as _slide_captions.json")
    args = parser.parse_args()

    video_path = os.path.abspath(args.video)
    if not os.path.exists(video_path):
        print(f"ERROR: Video file not found: {video_path}", file=sys.stderr)
        sys.exit(1)

    video_dir = os.path.dirname(video_path)
    output_path = args.output or os.path.join(video_dir, "_slide_terms.txt")

    print(f"Video: {video_path}", file=sys.stderr)
    print(f"Output: {output_path}", file=sys.stderr)

    ext = os.path.splitext(video_path)[1].lower()
    if ext in (".pptx", ".ppt"):
        if args.caption:
            print("pptx 無時間戳，忽略 --caption，只輸出 _slide_terms.txt", file=sys.stderr)

        lines = extract_pptx_text(video_path)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("# 本集投影片術語（pptx 抽取：標題/內文/表格/備註）\n")
            f.write(f"# 抽取時間: {time.strftime('%Y-%m-%d %H:%M')}\n")
            f.write("\n")
            f.write("\n".join(lines))

        print("\n=== Done (PowerPoint mode) ===", file=sys.stderr)
        print(f"  Slides: {video_path}", file=sys.stderr)
        print(f"  Extracted lines: {len(lines)}", file=sys.stderr)
        print(f"  Output: {output_path}", file=sys.stderr)
        return

    # Step 1: Extract frames
    with tempfile.TemporaryDirectory() as tmpdir:
        print(f"\n[1/4] Extracting frames (1 per {args.interval}s)...", file=sys.stderr)
        t0 = time.time()
        try:
            frames = extract_frames(video_path, tmpdir, args.interval)
        except RuntimeError as exc:
            print(f"ERROR: {exc}", file=sys.stderr)
            sys.exit(1)
        print(f"  → {len(frames)} frames extracted in {time.time()-t0:.1f}s", file=sys.stderr)

        # Step 2: Deduplicate
        print(f"\n[2/4] Deduplicating (threshold={args.threshold})...", file=sys.stderr)
        t0 = time.time()
        unique_frames = deduplicate_frames(frames, args.threshold)
        print(f"  → {len(unique_frames)} unique frames (removed {len(frames)-len(unique_frames)} duplicates) in {time.time()-t0:.1f}s", file=sys.stderr)

        if not unique_frames:
            print("  No unique frames found (video may not contain slides)", file=sys.stderr)
            write_terms_file({"tickers": [], "proper_nouns": [], "technical_terms": [], "slide_titles": []}, output_path)
            return

        # Step 3: OCR/Caption with VLM
        mode_label = "caption" if args.caption else "OCR"
        try:
            engine, model = resolve_engine(args)
        except RuntimeError as exc:
            print(f"ERROR: {exc}", file=sys.stderr)
            sys.exit(1)

        backend = {"ollama": "Ollama", "mlx": "mlx-vlm", "apple-vision": "Apple Vision", "rapidocr": "RapidOCR"}[engine]
        model_label = f" with {model}" if model else ""
        print(f"\n[3/4] Running {mode_label}{model_label} ({backend})...", file=sys.stderr)
        t0 = time.time()
        try:
            if engine == "rapidocr":
                raw_results = ocr_with_rapidocr(unique_frames)
            elif engine == "apple-vision":
                raw_results = ocr_with_apple_vision(unique_frames)
            else:
                raw_results = run_vlm_engine(engine, unique_frames, model, args.caption)
        except Exception as exc:
            if args.engine != "auto" or engine not in OCR_ENGINES:
                print(f"ERROR: {backend} failed: {exc}", file=sys.stderr)
                sys.exit(1)
            attempted = [backend]
            _warn(f"{backend} failed; falling back to Ollama {OLLAMA_DEFAULT_MODEL}: {exc}")
            engine, model, backend = "ollama", OLLAMA_DEFAULT_MODEL, "Ollama"
            try:
                raw_results = run_vlm_engine(engine, unique_frames, model, args.caption)
            except Exception as ollama_exc:
                attempted.append(backend)
                _warn(f"Ollama fallback failed; trying mlx-vlm {MLX_DEFAULT_MODEL}: {ollama_exc}")
                engine, model, backend = "mlx", MLX_DEFAULT_MODEL, "mlx-vlm"
                try:
                    raw_results = run_vlm_engine(engine, unique_frames, model, args.caption)
                except Exception as mlx_exc:
                    attempted.append(backend)
                    print(f"ERROR: all engines failed ({' → '.join(attempted)}); last error: {mlx_exc}", file=sys.stderr)
                    sys.exit(1)
        print(f"  → {mode_label} completed in {time.time()-t0:.1f}s", file=sys.stderr)

    # Step 4: Parse and write
    print(f"\n[4/4] Parsing and writing output...", file=sys.stderr)
    parsed = parse_ocr_outputs(raw_results, caption_mode=args.caption) if engine in OCR_ENGINES else parse_vlm_outputs(raw_results, caption_mode=args.caption)

    if args.caption:
        # Caption mode: write timestamped captions JSON
        captions_path = output_path.replace("_slide_terms.txt", "_slide_captions.json")
        if captions_path == output_path:
            captions_path = output_path.replace(".txt", "_captions.json")
        with open(captions_path, "w", encoding="utf-8") as f:
            json.dump(parsed["captions"], f, ensure_ascii=False, indent=2)

        # Also write a terms file for backward compatibility
        compat_terms = {
            "tickers": [], "proper_nouns": [],
            "technical_terms": parsed["all_terms"],
            "slide_titles": [],
        }
        write_terms_file(compat_terms, output_path)

        print(f"\n=== Done (caption mode) ===", file=sys.stderr)
        print(f"  Captions: {len(parsed['captions'])} timestamped entries", file=sys.stderr)
        print(f"  Unique terms: {len(parsed['all_terms'])}", file=sys.stderr)
        print(f"  Captions JSON: {captions_path}", file=sys.stderr)
        print(f"  Terms (compat): {output_path}", file=sys.stderr)
    else:
        # Legacy OCR mode
        write_terms_file(parsed, output_path)
        total = len(parsed["tickers"]) + len(parsed["proper_nouns"]) + len(parsed["technical_terms"])
        print(f"\n=== Done ===", file=sys.stderr)
        print(f"  Tickers: {len(parsed['tickers'])} — {', '.join(parsed['tickers']) or '(none)'}", file=sys.stderr)
        print(f"  Proper nouns: {len(parsed['proper_nouns'])}", file=sys.stderr)
        print(f"  Technical terms: {len(parsed['technical_terms'])}", file=sys.stderr)
        print(f"  Slide titles: {len(parsed['slide_titles'])}", file=sys.stderr)
        print(f"  Total unique terms: {total}", file=sys.stderr)
        print(f"  Output: {output_path}", file=sys.stderr)

    if args.json:
        json_path = output_path.replace(".txt", ".json")
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump({"parsed": parsed, "raw_results": [{k: v for k, v in r.items() if k != "frame"} for r in raw_results]},
                      f, ensure_ascii=False, indent=2)
        print(f"  JSON: {json_path}", file=sys.stderr)


if __name__ == "__main__":
    main()
